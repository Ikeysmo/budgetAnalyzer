# This is a sample Python script.

# Press Shift+F10 to execute it or replace it with your code.
# Press Double Shift to search everywhere for classes, files, tool windows, actions, and settings.

import pandas
import re
from rules_dict import rules_dict
import graphs
import statistics
import PPTWrapper as mypt
import matplotlib.ticker as tick

months_idx = ['Jan', 'Feb', 'March', 'April', 'May', 'June', 'July', "Aug", "Sept", "Oct", "Nov", "Dec"]

def __getCategory(title):
    for categories, list_dicts in rules_dict.items():
        for name, regex_match in list_dicts.items():
            if type(regex_match) not in [list, set] :
                regex_match = [regex_match]
            for sub_match in regex_match:
                if re.search(sub_match, title.lower()):
                    return categories

    return "NA"


def __getCommonName(title):
    for categories, list_dicts in rules_dict.items():
        for name, regex_match in list_dicts.items():
            if type(regex_match) not in [list, set] :
                regex_match = [regex_match]
            for sub_match in regex_match:
                if re.search(sub_match, title.lower()):
                    return name

    return "NA"


def parseSpreadsheet(filepath = r'ufcu_joint_budget_pull_History-020523-052449.csv'):
    budg_table = pandas.read_csv(filepath,
                                 names=["Date", "Title", "Cost", "Running Amt"])
    print(budg_table.head())
    budg_table[budg_table.columns[2:]] = budg_table[budg_table.columns[2:]].replace('[\$,]', '', regex=True).astype(
        float)
    budg_table["Category"] = budg_table["Title"].map(__getCategory)
    budg_table["Common Name"] = budg_table["Title"].map(__getCommonName)
    parsed_date_table = budg_table["Date"].str.extract("(?P<Month>\d+)/(?P<Day>\d+)/(?P<Year>\d+)")

    budg_table = pandas.concat([budg_table, parsed_date_table], axis=1)
    budg_table['Month'] = pandas.to_numeric(budg_table['Month'])
    budg_table['Day'] = pandas.to_numeric(budg_table['Day'])
    budg_table['Year'] = pandas.to_numeric(budg_table['Year'])

    NA_remaining = budg_table[budg_table["Category"] == "NA"]

    income_or_refunds = budg_table[budg_table["Cost"] > 0]
    expenses = budg_table[budg_table["Cost"] < 0]

    new_excel_dict = {"Total" : budg_table, "Income": income_or_refunds}

    categories = set(expenses["Category"].values)
    # expenses.to_csv("ufcu_joint_budget_parsed_output.csv", index=False)

    for cate in categories:
        new_excel_dict[cate] = expenses[expenses["Category"] == cate]

    return new_excel_dict

# Press the green button in the gutter to run the script.
if __name__ == '__main__':
    from pptx import Presentation
    from pptx.util import Inches
    import os
    import matplotlib.pyplot as plt

    prs = Presentation()


    # food_only = expenses[expenses["Category"] == "FOOD"]
    # food_only.to_csv("xxx_ufcu_joint_budget_parsed_output.csv", index=False)
    new_excel_dict  = parseSpreadsheet()

    mypt.add_title_slide(prs, "Budget Analysis", "Study of budget")

    mplt = plt.figure(figsize=(10, 5))
    print(type(mplt))
    graphs.graphBarCategoryTotals(new_excel_dict['Total'], mplt)

    mplt2 = plt.figure(figsize=(10, 5))
    graphs.graphBarMultiple(new_excel_dict['Total'], mplt2)



    mypt.add_picture_slide(prs, mplt, "Monthly")
    mypt.add_picture_slide(prs, mplt2, "Monthly")


    mcategories = set(new_excel_dict['Total']['Category'])
    import seaborn

    for _cate in mcategories:
        tf = statistics.summ_category_by_month_only(new_excel_dict['Total'], _cate)
        plt.figure(figsize=(10,5))
        #aw = seaborn.lineplot(tf, x=tf['Month'], y=tf['Cost'], marker='o', )
        aw = seaborn.barplot(tf, x=tf['Month'], y=tf['Cost'] )
        for p in aw.patches:
            height = p.get_height()
            aw.annotate('${:.2f}'.format(height),
                        xy=(p.get_x() + p.get_width() / 2, height),
                        xytext=(0, 3),  # 3 points vertical offset
                        textcoords="offset points",
                        ha='center', va='bottom')
        aw.set_title("Category {} vs Month".format(_cate))
        aw.grid(b=True)
        aw.yaxis.set_major_formatter(tick.StrMethodFormatter('${x}'))
        #plt.show()
        mypt.add_picture_with_content_slide(prs, plt, "Expenses by Month for Category {}".format(_cate))
    #plt.show()

    months = set(new_excel_dict['Total']['Month'])

    for _m in months:
        mplt = plt.figure(figsize=(10, 9))
        tf = statistics.summ_category_by_month(new_excel_dict['Total'], _m, exclude_columns=["TRAVEL"])
        graphs.graphPieChart(tf, "Categories vs Month - {}".format(months_idx[_m-1]))
        mypt.add_picture_with_content_slide(prs, plt, "Total Money Spent for Month {} (Percentage Wise)".format(months_idx[_m-1]))

    # tf = statistics.summ_category_by_month(new_excel_dict['Total'],7, exclude_columns=["TRAVEL"])
    # graphs.graphPieChart(tf, "Categories vs Month(2)")
    #graphs.plt.show()

    prs.save('c:\\temp\\myoutput2.pptx')
    #graphs.graphThis(new_excel_dict['Total'])

    # with pandas.ExcelWriter("myExcelBudgetSummary.xlsx") as writer:
        # for k,v in new_excel_dict.items():
        #     v.to_excel(writer, sheet_name = k, index=False)

    os.startfile('c:\\temp\\myoutput2.pptx')


