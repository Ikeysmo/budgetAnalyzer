
import matplotlib.pyplot as plt
import matplotlib.ticker as tick
import PPTWrapper as mypt
import seaborn

def graphBarCategoryTotals(pandas_table, fig = None):

    categories = set(pandas_table["Category"])

    #filter out work
    pandas_table = pandas_table[pandas_table["Category"] != "WORK"]
    pandas_table = pandas_table[pandas_table["Category"] != "Bank"]

    #make numbers positive
    pandas_table['Cost'] = pandas_table['Cost'].abs()


    summ_table = pandas_table.groupby(['Category']).sum(numeric_only=True)
    if not fig:
        fig = plt.figure(figsize=(10,5))
    ax = fig.add_subplot()
    ax.yaxis.set_major_formatter(tick.StrMethodFormatter('${x}'))

    # plt.subplot(111, autoscale_on=True)
    pps = ax.bar(summ_table.index, summ_table['Cost'])
    ax.grid(linestyle='--' )
    ax.set_title("Costs Tracker")
    plt.xlabel('Category')
    plt.ylabel("Cost")

    for p in pps:
        height = p.get_height()
        ax.annotate('${}'.format(height),
                    xy=(p.get_x() + p.get_width() / 2, height),
                    xytext=(0, 3),  # 3 points vertical offset
                    textcoords="offset points",
                    ha='center', va='bottom')



    summ_table2 = pandas_table["Cost"].agg(['sum'])
    print("Done")



def graphBarMultiple(pandas_table, fig=None):
    categories = set(pandas_table["Category"])

    #filter out work
    pandas_table = pandas_table[pandas_table["Category"] != "WORK"]
    pandas_table = pandas_table[pandas_table["Category"] != "Bank"]
    pandas_table = pandas_table[pandas_table["Category"] != "HOUSE"]

    #make numbers positive
    pandas_table['Cost'] = pandas_table['Cost'].abs()


    summ_table = pandas_table.groupby(['Category', 'Month'], as_index=False).sum(numeric_only=True)

    if not fig:
        fig = plt.figure(figsize=(10,5))
    ax = fig.add_subplot()
    ax.yaxis.set_major_formatter(tick.StrMethodFormatter('${x}'))


    aw = seaborn.barplot(data=summ_table, x="Month", y="Cost", hue='Category', width = 1, order=sorted(set(summ_table["Month"])))
    aw.grid(b=True)
    aw.set_title("Categories vs Month")
    print("Done")


def graphBarSingleExpense(pandas_table, expense_category = None, fig=None):
    categories = set(pandas_table["Category"])

    #filter out work
    pandas_table = pandas_table[pandas_table["Category"] ==expense_category]

    #make numbers positive
    pandas_table['Cost'] = pandas_table['Cost'].abs()


    summ_table = pandas_table.groupby(['Category', 'Month'], as_index=False).sum(numeric_only=True)

    if not fig:
        fig = plt.figure(figsize=(10,5))
    ax = fig.add_subplot()
    ax.yaxis.set_major_formatter(tick.StrMethodFormatter('${x}'))


    aw = seaborn.barplot(data=summ_table, x="Month", y="Cost",  order=sorted(set(summ_table["Month"])))
    aw.grid(b=True)
    aw.set_title("Category ({}) vs Month".format(expense_category))
    print("Done")


def graphBarMultiple(pandas_table, fig=None):
    categories = set(pandas_table["Category"])

    #filter out work
    pandas_table = pandas_table[pandas_table["Category"] != "WORK"]
    pandas_table = pandas_table[pandas_table["Category"] != "Bank"]
    pandas_table = pandas_table[pandas_table["Category"] != "HOUSE"]

    #make numbers positive
    pandas_table['Cost'] = pandas_table['Cost'].abs()


    summ_table = pandas_table.groupby(['Category', 'Month'], as_index=False).sum(numeric_only=True)

    if not fig:
        fig = plt.figure(figsize=(10,5))
    ax = fig.add_subplot()
    ax.yaxis.set_major_formatter(tick.StrMethodFormatter('${x}'))


    aw = seaborn.barplot(data=summ_table, x="Month", y="Cost", hue='Category', width = 1, order=sorted(set(summ_table["Month"])))
    aw.grid(b=True)
    aw.set_title("Categories vs Month")
    print("Done")





def graphThis(pandas_table):
    from PIL import Image
    from pptx import Presentation
    from pptx.util import Inches

    from datetime import date


    mplt = plt.figure(figsize=(10,5))
    print(type(mplt))
    graphBarCategoryTotals(pandas_table, mplt)

    mplt2 = plt.figure(figsize=(10,5))
    graphBarMultiple(pandas_table, mplt2)

    mplt3 = plt.figure(figsize=(10,5))
    graphBarSingleExpense(pandas_table, "FOOD", mplt3)
    plt.show()


    mplt4 = plt.figure(figsize=(10,5))
    graphBarSingleExpense(pandas_table, "FOOD", mplt4)
    plt.show()




    if False:
        prs = Presentation()


        mypt.add_title_slide(prs, "Budget Analysis", "Study of budget")
        mypt.add_picture_slide(prs, mplt, "Monthly")
        mypt.add_picture_slide(prs, mplt2, "Monthly")
        mypt.add_picture_slide(prs, mplt3, "Monthly")

        prs.save('c:\\temp\\myoutput2.pptx')
        import os
        os.startfile('c:\\temp\\myoutput2.pptx')
